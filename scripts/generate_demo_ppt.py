#!/usr/bin/env python3
"""Generate a polished demo presentation for Campus Course Review Explorer.

Fully custom layouts (blank slides + manual shapes) for a modern, dark theme
that mirrors the app UI. No default placeholder bullets.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Campus_Course_Review_Demo.pptx"

# ----- Palette (matches the dark slate / indigo app theme) -----
BG_DARK = RGBColor(0x0B, 0x11, 0x20)      # deep slate
BG_PANEL = RGBColor(0x15, 0x1E, 0x33)     # card slate
BG_PANEL_2 = RGBColor(0x1C, 0x27, 0x42)   # lighter card
BRAND = RGBColor(0x6366F1 >> 16, (0x6366F1 >> 8) & 0xFF, 0x6366F1 & 0xFF)  # indigo
BRAND_LT = RGBColor(0x81, 0x8C, 0xF8)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)       # cyan
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)

SW = Inches(13.333)
SH = Inches(7.5)

SPEAKER_COLORS = [BRAND, ACCENT, GREEN, AMBER, BRAND_LT]


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def _rounded(slide, x, y, w, h, color, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(shape, color)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.shadow.inherit = False
    return shape


def _rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(shape, color)
    shape.shadow.inherit = False
    return shape


def _text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
    spacing=1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def _bullets(slide, x, y, w, h, items, size=16, color=WHITE, gap=10, marker_color=ACCENT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        marker = p.add_run()
        marker.text = "▸  "
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = marker_color
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def _eyebrow(slide, x, y, text, color=ACCENT):
    box = _text(slide, x, y, Inches(6), Inches(0.4), text.upper(), size=13, color=color, bold=True)
    box.text_frame.paragraphs[0].runs[0].font._rPr.set("spc", "200")
    return box


def _footer(slide, idx):
    _text(
        slide,
        Inches(0.7),
        Inches(7.0),
        Inches(8),
        Inches(0.3),
        "Campus Course Review · Sentiment & Topic Explorer",
        size=10,
        color=DIM,
    )
    _text(
        slide,
        Inches(12.0),
        Inches(7.0),
        Inches(0.8),
        Inches(0.3),
        str(idx),
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    return slide


def content_header(slide, eyebrow, title, accent=ACCENT):
    _rect(slide, Inches(0.7), Inches(0.6), Inches(0.12), Inches(0.95), accent)
    _eyebrow(slide, Inches(0.95), Inches(0.6), eyebrow, accent)
    _text(slide, Inches(0.95), Inches(0.92), Inches(11.5), Inches(0.8), title, size=30, bold=True)


# ----------------------------------------------------------------------------
def slide_title(prs):
    slide = blank(prs)
    # accent band
    _rect(slide, 0, 0, SW, Inches(0.18), BRAND)
    _eyebrow(slide, Inches(1.0), Inches(1.7), "NLP & Machine Learning · Capstone Demo", ACCENT)
    _text(
        slide,
        Inches(1.0),
        Inches(2.2),
        Inches(11.3),
        Inches(2.0),
        "Campus Course Review\nSentiment & Topic Explorer",
        size=52,
        bold=True,
        spacing=1.0,
    )
    _text(
        slide,
        Inches(1.0),
        Inches(4.3),
        Inches(11),
        Inches(0.6),
        "Turning thousands of unstructured course reviews into clear, actionable insight.",
        size=19,
        color=MUTED,
    )
    # stat chips
    chips = [("46", "Universities"), ("~1,000", "Reviews"), ("4", "AI summary engines"), ("100%", "Local MVP")]
    cx = Inches(1.0)
    for value, label in chips:
        card = _rounded(slide, cx, Inches(5.2), Inches(2.7), Inches(1.1), BG_PANEL, radius=0.12)
        _text(slide, cx, Inches(5.32), Inches(2.7), Inches(0.5), value, size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _text(slide, cx, Inches(5.85), Inches(2.7), Inches(0.4), label, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        cx += Inches(2.95)
    _notes(slide, "Intro the project: a dashboard that summarizes course reviews using NLP and ML. "
                  "Built as a working local MVP with real public data across 46 universities.")
    return slide


def slide_section(prs, number, title, speaker, color):
    slide = blank(prs)
    _set_bg(slide, BG_PANEL)
    _rect(slide, 0, 0, Inches(0.35), SH, color)
    _text(slide, Inches(1.0), Inches(2.2), Inches(3), Inches(1.6), f"{number:02d}", size=96, bold=True, color=color)
    _text(slide, Inches(1.05), Inches(3.7), Inches(11), Inches(1.0), title, size=40, bold=True)
    _rounded(slide, Inches(1.05), Inches(4.9), Inches(3.4), Inches(0.6), BG_PANEL_2, radius=0.5)
    _text(slide, Inches(1.05), Inches(4.96), Inches(3.4), Inches(0.48), f"🎤  {speaker}", size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_problem(prs):
    slide = blank(prs)
    content_header(slide, "Person 1 · Problem Statement", "Reading hundreds of reviews doesn't scale", RED)
    _text(slide, Inches(0.95), Inches(1.9), Inches(11.4), Inches(0.7),
          "\u201cStudents rely on large volumes of course reviews, but reading them manually is "
          "slow and makes trends hard to spot.\u201d", size=17, color=MUTED, italic=True)
    pains = [
        ("⏳", "Time-consuming", "Hundreds of reviews per course — hours of manual reading."),
        ("🧩", "Hidden patterns", "Workload, grading, and exam themes are buried in free text."),
        ("⚖️", "Hard to compare", "No easy way to weigh sentiment across courses or professors."),
        ("🎯", "Poor decisions", "Students pick courses on incomplete, anecdotal information."),
    ]
    x = Inches(0.95)
    y = Inches(2.9)
    w = Inches(5.55)
    h = Inches(1.7)
    for i, (icon, head, body) in enumerate(pains):
        col = i % 2
        row = i // 2
        cx = x + col * Inches(5.85)
        cy = y + row * Inches(1.95)
        _rounded(slide, cx, cy, w, h, BG_PANEL, radius=0.08)
        _text(slide, cx + Inches(0.3), cy + Inches(0.25), Inches(1), Inches(0.7), icon, size=30)
        _text(slide, cx + Inches(1.2), cy + Inches(0.28), w - Inches(1.4), Inches(0.5), head, size=18, bold=True, color=WHITE)
        _text(slide, cx + Inches(1.2), cy + Inches(0.78), w - Inches(1.4), Inches(0.8), body, size=13, color=MUTED)
    _notes(slide, "Set up the pain: manual review reading is slow and patterns are hard to see. "
                  "Our goal is to use NLP & ML to surface insight automatically.")
    _footer(slide, 3)
    return slide


def slide_solution(prs):
    slide = blank(prs)
    content_header(slide, "Person 2 · Solution Overview", "An insight dashboard for course reviews", ACCENT)
    _text(slide, Inches(0.95), Inches(1.95), Inches(11.4), Inches(0.8),
          "Browse by university → course → professor / semester / sentiment, and get a summarized view in seconds.",
          size=16, color=MUTED)
    features = [
        ("📊", "Sentiment Analysis", "Positive / neutral / negative per review."),
        ("🏷️", "Topic Extraction", "Workload, Exams, Grading, Lectures, Projects."),
        ("📈", "Insights Dashboard", "KPIs, charts, and an AI-written summary."),
        ("🔍", "Review Explorer", "Search & filter individual student reviews."),
        ("🤖", "AI Summaries", "Choose OpenAI, Groq, or Ollama at runtime."),
        ("🔐", "Secure Submissions", "Firebase-authenticated review posting."),
    ]
    x0 = Inches(0.95)
    y0 = Inches(2.85)
    w = Inches(3.7)
    h = Inches(1.75)
    for i, (icon, head, body) in enumerate(features):
        col = i % 3
        row = i // 3
        cx = x0 + col * Inches(3.85)
        cy = y0 + row * Inches(1.95)
        _rounded(slide, cx, cy, w, h, BG_PANEL, radius=0.09)
        _rect(slide, cx, cy, Inches(0.1), h, BRAND)
        _text(slide, cx + Inches(0.3), cy + Inches(0.22), Inches(1), Inches(0.6), icon, size=26)
        _text(slide, cx + Inches(0.3), cy + Inches(0.85), w - Inches(0.5), Inches(0.45), head, size=15, bold=True)
        _text(slide, cx + Inches(0.3), cy + Inches(1.25), w - Inches(0.5), Inches(0.5), body, size=11.5, color=MUTED)
    _notes(slide, "Pitch the product and walk through the six core features. Emphasize the goal: "
                  "help students make informed decisions faster.")
    _footer(slide, 5)
    return slide


def slide_architecture(prs):
    slide = blank(prs)
    content_header(slide, "Person 3 · Architecture & Design", "How data flows through the system", GREEN)
    steps = [
        ("Course Reviews", "RMP public sample + live user submissions", BRAND),
        ("Data Processing", "Import, clean, normalize → SQLite", ACCENT),
        ("NLP & Sentiment", "Keyword model (bulk) + LLM (live)", GREEN),
        ("Topic Extraction", "RMP tags + AI / keyword labels", AMBER),
        ("Interactive Dashboard", "React UI ← FastAPI REST API", BRAND_LT),
    ]
    y = Inches(2.3)
    x = Inches(0.95)
    w = Inches(2.18)
    h = Inches(2.2)
    gap = Inches(0.18)
    for i, (head, body, color) in enumerate(steps):
        cx = x + i * (w + gap)
        _rounded(slide, cx, y, w, h, BG_PANEL, radius=0.08)
        _rect(slide, cx, y, w, Inches(0.12), color)
        _text(slide, cx, y + Inches(0.35), w, Inches(0.5), f"{i+1}", size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        _text(slide, cx + Inches(0.18), y + Inches(0.95), w - Inches(0.36), Inches(0.6), head, size=14, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, cx + Inches(0.18), y + Inches(1.5), w - Inches(0.36), Inches(0.6), body, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + w + Emu(1), y + Inches(0.9), gap, Inches(0.4))
            _solid(arrow, DIM)
    _text(slide, Inches(0.95), Inches(4.9), Inches(11.4), Inches(0.5),
          "Current phase: runs fully locally on real public data — no cloud dependency required.",
          size=14, color=MUTED, italic=True)
    # tech stack strip
    _text(slide, Inches(0.95), Inches(5.55), Inches(11), Inches(0.4), "TECH STACK", size=12, bold=True, color=ACCENT)
    stack = ["React + Vite + Tailwind", "FastAPI + SQLAlchemy", "SQLite / Postgres", "Firebase Auth", "OpenAI · Groq · Ollama"]
    cx = Inches(0.95)
    for item in stack:
        chip_w = Inches(0.18 + 0.105 * len(item))
        _rounded(slide, cx, Inches(5.95), chip_w, Inches(0.55), BG_PANEL_2, radius=0.5)
        _text(slide, cx, Inches(6.04), chip_w, Inches(0.4), item, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        cx += chip_w + Inches(0.2)
    _notes(slide, "Walk the pipeline left to right. Stress that it is a clean, layered architecture "
                  "that already runs locally and is ready to move to the cloud.")
    _footer(slide, 7)
    return slide


def slide_ml(prs):
    slide = blank(prs)
    content_header(slide, "Person 4 · Machine Learning & Analytics", "From raw text to classified insight", AMBER)
    # left: how it works
    _rounded(slide, Inches(0.95), Inches(1.95), Inches(6.0), Inches(4.4), BG_PANEL, radius=0.06)
    _text(slide, Inches(1.25), Inches(2.15), Inches(5.4), Inches(0.4), "How the ML works", size=18, bold=True)
    _bullets(slide, Inches(1.25), Inches(2.75), Inches(5.4), Inches(3.4), [
        "Bulk import: fast keyword sentiment + RMP topic tags",
        "Live reviews: one AI call → sentiment + topics together",
        "Summaries: selectable engine (OpenAI / Groq / Ollama / template)",
        "Dashboard aggregates topic counts and sentiment %",
        "Cross-university topic analytics (BigQuery-ready)",
    ], size=14, gap=14)
    # right: sentiment legend + example
    _rounded(slide, Inches(7.2), Inches(1.95), Inches(5.15), Inches(2.0), BG_PANEL, radius=0.06)
    _text(slide, Inches(7.45), Inches(2.12), Inches(4.6), Inches(0.4), "Sentiment classes", size=16, bold=True)
    legend = [("Positive", GREEN), ("Neutral", AMBER), ("Negative", RED)]
    lx = Inches(7.45)
    for label, color in legend:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx, Inches(2.75), Inches(0.28), Inches(0.28))
        _solid(dot, color)
        _text(slide, lx + Inches(0.38), Inches(2.71), Inches(1.3), Inches(0.4), label, size=13, color=WHITE)
        lx += Inches(1.6)
    _text(slide, Inches(7.45), Inches(3.25), Inches(4.7), Inches(0.6),
          "Each review → label + confidence + up to 3 topics", size=12, color=MUTED)
    # example callout
    _rounded(slide, Inches(7.2), Inches(4.1), Inches(5.15), Inches(2.25), BG_PANEL_2, radius=0.06)
    _text(slide, Inches(7.45), Inches(4.28), Inches(4.7), Inches(0.4), "💡  Example insight", size=15, bold=True, color=ACCENT)
    _text(slide, Inches(7.45), Inches(4.78), Inches(4.7), Inches(1.5),
          "Instead of reading 200 reviews, a student instantly sees:\n"
          "“Highly rated, but heavy workload and project-intensive.”",
          size=13.5, color=WHITE, spacing=1.1)
    _notes(slide, "Explain sentiment + topic classification and confidence. Use the example: a student "
                  "gets the gist in seconds instead of reading 200 reviews.")
    _footer(slide, 9)
    return slide


def slide_demo(prs):
    slide = blank(prs)
    content_header(slide, "Person 5 · Live Demo", "Five minutes, end to end", BRAND_LT)
    steps = [
        "Search a university (e.g. UIUC, Brooklyn College)",
        "Open a course dashboard (e.g. ASTR)",
        "Review KPIs, sentiment pie & topic distribution",
        "Switch AI summary model (OpenAI / Groq) → regenerate",
        "Filter reviews by sentiment / professor / semester",
        "Sign in → submit a review → watch AI analyze it live",
    ]
    y = Inches(2.05)
    for i, step in enumerate(steps):
        cy = y + i * Inches(0.72)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), cy, Inches(0.55), Inches(0.55))
        _solid(circle, BRAND)
        _text(slide, Inches(0.95), cy + Inches(0.02), Inches(0.55), Inches(0.5), str(i + 1), size=18, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _rounded(slide, Inches(1.7), cy, Inches(7.2), Inches(0.55), BG_PANEL, radius=0.2)
        _text(slide, Inches(1.95), cy + Inches(0.02), Inches(6.8), Inches(0.5), step, size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # URLs panel
    _rounded(slide, Inches(9.2), Inches(2.05), Inches(3.15), Inches(2.4), BG_PANEL_2, radius=0.08)
    _text(slide, Inches(9.45), Inches(2.25), Inches(2.7), Inches(0.4), "RUN LOCALLY", size=12, bold=True, color=ACCENT)
    _text(slide, Inches(9.45), Inches(2.8), Inches(2.7), Inches(1.5),
          "Frontend\n:5174\n\nBackend API\n:8080", size=14, color=WHITE, spacing=1.1)
    _notes(slide, "Drive the live demo with this checklist. Have the app open at localhost:5174 and the "
                  "backend running on :8080 beforehand.")
    _footer(slide, 11)
    return slide


def slide_roadmap(prs):
    slide = blank(prs)
    content_header(slide, "Roadmap", "Built today, scaling tomorrow", GREEN)
    cols = [
        ("✅ In the MVP", GREEN, [
            "Multi-university (46 schools)",
            "Sentiment + topic extraction",
            "LLM summaries (OpenAI/Groq/Ollama)",
            "Cross-university analytics",
        ]),
        ("Phase 2 · Cloud", ACCENT, [
            "Deploy to GCP / AWS",
            "Public web access",
            "Vertex AI / Gemini option",
            "Managed Postgres",
        ]),
        ("Phase 3 · Scale", AMBER, [
            "Larger datasets",
            "Scheduled re-imports",
            "Performance tuning",
            "Caching layer",
        ]),
        ("Phase 4 · Intelligence", BRAND_LT, [
            "Course recommender",
            "Professor comparison",
            "Trend forecasting",
            "Mobile app",
        ]),
    ]
    x0 = Inches(0.95)
    w = Inches(2.85)
    h = Inches(4.0)
    for i, (head, color, items) in enumerate(cols):
        cx = x0 + i * Inches(2.98)
        _rounded(slide, cx, Inches(2.1), w, h, BG_PANEL, radius=0.06)
        _rect(slide, cx, Inches(2.1), w, Inches(0.7), color)
        _text(slide, cx + Inches(0.2), Inches(2.27), w - Inches(0.4), Inches(0.45), head, size=14, bold=True, color=BG_DARK, anchor=MSO_ANCHOR.MIDDLE)
        _bullets(slide, cx + Inches(0.25), Inches(3.0), w - Inches(0.45), Inches(2.9), items, size=12.5, gap=11, marker_color=color)
    _notes(slide, "Be honest about what's already done (left column) versus future phases. "
                  "Multi-university and LLM summaries are already working, not just planned.")
    _footer(slide, 12)
    return slide


def slide_closing(prs):
    slide = blank(prs)
    _set_bg(slide, BG_PANEL)
    _rect(slide, 0, 0, SW, Inches(0.18), BRAND)
    _eyebrow(slide, Inches(1.0), Inches(1.9), "Closing", ACCENT)
    _text(slide, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.0),
          "From unstructured feedback\nto meaningful insight.", size=44, bold=True, spacing=1.0)
    _text(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(1.0),
          "A locally hosted prototype with a production-ready, cloud-scalable architecture — "
          "showing how ML & NLP turn student feedback into better academic decisions.",
          size=17, color=MUTED, spacing=1.15)
    _rounded(slide, Inches(1.0), Inches(5.9), Inches(3.0), Inches(0.7), BRAND, radius=0.3)
    _text(slide, Inches(1.0), Inches(6.0), Inches(3.0), Inches(0.5), "Thank you — Questions?", size=15, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _notes(slide, "Wrap up: reinforce the value proposition and that the architecture is ready to scale "
                  "to the cloud in the next phase. Invite questions.")
    return slide


def main() -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_title(prs)
    slide_section(prs, 1, "Problem Statement", "Person 1", RED)
    slide_problem(prs)
    slide_section(prs, 2, "Solution Overview", "Person 2", ACCENT)
    slide_solution(prs)
    slide_section(prs, 3, "Architecture & Design", "Person 3", GREEN)
    slide_architecture(prs)
    slide_section(prs, 4, "Machine Learning & Analytics", "Person 4", AMBER)
    slide_ml(prs)
    slide_section(prs, 5, "Live Demo & Roadmap", "Person 5", BRAND_LT)
    slide_demo(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
